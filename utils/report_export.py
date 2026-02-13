"""
Export competitor intelligence report to Word (.docx) and PowerPoint (.pptx)
with logo and brand theme (#d03e9d, #08216b).
"""
import io
import re
from pathlib import Path
from typing import Optional, List, Tuple

# Theme colors (RGB 0-255)
PRIMARY_PINK = (208, 62, 157)   # #d03e9d
PRIMARY_NAVY = (8, 33, 107)     # #08216b
SUBTLE_BG = (245, 247, 251)    # #f5f7fb


def _svg_to_png_bytes(svg_path: str) -> Optional[bytes]:
    """Convert SVG file to PNG bytes for embedding in Word/PPT. Returns None if conversion fails.
    Requires cairo system library; if missing (e.g. on Windows), returns None so export still works without logo."""
    try:
        import cairosvg
    except (ImportError, OSError, Exception):
        return None
    try:
        with open(svg_path, "rb") as f:
            png_bytes = cairosvg.svg2png(file_obj=f, output_width=280, dpi=150)
        return png_bytes
    except Exception:
        return None


def _parse_report_blocks(report_text: str) -> List[Tuple[str, str]]:
    """Parse markdown-style report into (block_type, content) list. Types: 'heading', 'paragraph', 'bullet'."""
    blocks: List[Tuple[str, str]] = []
    lines = report_text.replace("\r\n", "\n").split("\n")
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        # ## Section heading
        if stripped.startswith("## "):
            blocks.append(("heading", stripped[3:].strip()))
            i += 1
            continue
        # # Heading (treat as heading level 2)
        if stripped.startswith("# ") and not stripped.startswith("## "):
            blocks.append(("heading", stripped[2:].strip()))
            i += 1
            continue
        # Bullet: - or * at start
        if stripped.startswith("- ") or stripped.startswith("* "):
            blocks.append(("bullet", stripped[2:].strip()))
            i += 1
            continue
        # Numbered line (e.g. "1. British Gas snapshot")
        if re.match(r"^\d+\.\s", stripped):
            blocks.append(("bullet", stripped))
            i += 1
            continue
        # Empty line
        if not stripped:
            i += 1
            continue
        # Paragraph: collect consecutive non-bullet/non-heading lines
        para_lines = [stripped]
        i += 1
        while i < len(lines):
            next_line = lines[i]
            next_stripped = next_line.strip()
            if not next_stripped or next_stripped.startswith("## ") or next_stripped.startswith("# ") or next_stripped.startswith("- ") or next_stripped.startswith("* ") or re.match(r"^\d+\.\s", next_stripped):
                break
            para_lines.append(next_stripped)
            i += 1
        blocks.append(("paragraph", " ".join(para_lines)))
    return blocks


def export_report_to_docx(report_text: str, logo_path: Optional[str] = None, title: str = "AI Powered Competitor Intelligence") -> bytes:
    """Generate a Word document from the report text. Optional logo at top. Returns .docx bytes."""
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    # Optional logo at top
    if logo_path and Path(logo_path).exists():
        if logo_path.lower().endswith(".svg"):
            png_bytes = _svg_to_png_bytes(logo_path)
            if png_bytes:
                try:
                    doc.add_picture(io.BytesIO(png_bytes), width=Inches(1.75))
                    last = doc.paragraphs[-1]
                    last.alignment = WD_ALIGN_PARAGRAPH.LEFT
                except Exception:
                    pass
        else:
            try:
                doc.add_picture(logo_path, width=Inches(1.75))
            except Exception:
                pass

    # Title
    p = doc.add_paragraph()
    p_run = p.add_run(title)
    p_run.bold = True
    p_run.font.size = Pt(18)
    p_run.font.color.rgb = _rgb_to_docx(PRIMARY_NAVY)
    p.paragraph_format.space_after = Pt(6)

    blocks = _parse_report_blocks(report_text)
    for block_type, content in blocks:
        if block_type == "heading":
            h = doc.add_heading(content, level=1)
            for run in h.runs:
                run.font.color.rgb = _rgb_to_docx(PRIMARY_NAVY)
                run.font.size = Pt(14)
        elif block_type == "bullet":
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(content)
        elif block_type == "paragraph" and content:
            doc.add_paragraph(content)

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _rgb_to_docx(rgb: Tuple[int, int, int]):
    from docx.shared import RGBColor
    return RGBColor(*rgb)


def _rgb_to_pptx(rgb: Tuple[int, int, int]):
    from pptx.dml.color import RGBColor
    return RGBColor(*rgb)


def _condense_for_slide(line: str, max_len: int = 72) -> str:
    """Shorten text for slide bullets; keep one line where possible."""
    line = line.strip()
    if not line:
        return ""
    # Remove leading bullet/asterisk
    for prefix in ("• ", "- ", "* ", "– "):
        if line.startswith(prefix):
            line = line[len(prefix):].strip()
            break
    if len(line) <= max_len:
        return line
    # Prefer break at last space before max_len
    cut = line[: max_len + 1].rfind(" ")
    if cut > max_len // 2:
        return line[:cut].strip() + "…"
    return line[:max_len].strip() + "…"


def _content_to_slide_bullets(content_list: List[str], max_bullets_per_slide: int = 6) -> List[List[str]]:
    """Split content into slide-sized chunks of condensed bullets (for slide-specific formatting)."""
    bullets: List[str] = []
    for line in content_list:
        line = line.strip()
        for prefix in ("• ", "- ", "* "):
            if line.startswith(prefix):
                line = line[len(prefix):].strip()
                break
        if not line:
            continue
        # One long paragraph → split into short bullets by sentence or length
        if len(line) > 120:
            parts = re.split(r"[.;]\s+", line)
            for part in parts:
                part = part.strip()
                if part:
                    bullets.append(_condense_for_slide(part, max_len=78))
        else:
            bullets.append(_condense_for_slide(line, max_len=78))

    # Chunk into slides
    result: List[List[str]] = []
    for i in range(0, len(bullets), max_bullets_per_slide):
        result.append(bullets[i : i + max_bullets_per_slide])
    return result if result else [[]]


def export_report_to_pptx(report_text: str, logo_path: Optional[str] = None, title: str = "AI Powered Competitor Intelligence") -> bytes:
    """Generate a formatted PowerPoint deck: theme colors, condensed slide-specific content, no blank slides."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    navy = RGBColor(*PRIMARY_NAVY)
    pink = RGBColor(*PRIMARY_PINK)
    white = RGBColor(255, 255, 255)
    body_gray = RGBColor(55, 65, 81)

    # Parse report into sections (title + content lines)
    blocks = _parse_report_blocks(report_text)
    sections: List[Tuple[str, List[str]]] = []
    current_title = "Competitor Intelligence"
    current_content: List[str] = []

    for block_type, content in blocks:
        if block_type == "heading":
            if current_content or current_title != "Competitor Intelligence":
                sections.append((current_title, current_content))
            current_title = content
            current_content = []
        elif block_type == "bullet":
            current_content.append(content)
        elif block_type == "paragraph" and content:
            current_content.append(content)

    if current_title or current_content:
        sections.append((current_title, current_content))

    # Logo
    logo_bytes: Optional[bytes] = None
    if logo_path and Path(logo_path).exists():
        if logo_path.lower().endswith(".svg"):
            logo_bytes = _svg_to_png_bytes(logo_path)
        else:
            try:
                with open(logo_path, "rb") as f:
                    logo_bytes = f.read()
            except Exception:
                pass

    def add_logo(slide, left: float, top: float, width: float = 1.25):
        if not logo_bytes:
            return
        try:
            slide.shapes.add_picture(io.BytesIO(logo_bytes), Inches(left), Inches(top), width=Inches(width))
        except Exception:
            return

    blank_layout = prs.slide_layouts[6]

    # ---------- Title slide: themed ----------
    slide = prs.slides.add_slide(blank_layout)
    # Navy bar full width at top
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.6)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = navy
    bar.line.fill.background()
    add_logo(slide, 0.6, 0.35)
    # Title text on bar (white) – right of logo
    title_box = slide.shapes.add_textbox(Inches(2.2), Inches(0.35), Inches(9), Inches(0.95))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = white
    # Subtitle below bar (pink)
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12), Inches(0.6))
    sub_box.text_frame.text = "Competitor Intelligence Report"
    sub_box.text_frame.paragraphs[0].font.size = Pt(18)
    sub_box.text_frame.paragraphs[0].font.color.rgb = pink
    # Optional footer line
    line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.2), Inches(2), Pt(4))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = pink
    line_shape.line.fill.background()

    # ---------- Content slides: navy title bar + condensed bullets ----------
    for section_title, content_list in sections:
        slide_chunks = _content_to_slide_bullets(content_list, max_bullets_per_slide=6)
        if not slide_chunks:
            slide_chunks = [[]]  # at least one slide with just the title

        for chunk_idx, bullets in enumerate(slide_chunks):
            slide = prs.slides.add_slide(blank_layout)

            # Navy title bar
            bar_height = Inches(0.95)
            bar_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, bar_height
            )
            bar_shape.fill.solid()
            bar_shape.fill.fore_color.rgb = navy
            bar_shape.line.fill.background()

            # Section title on bar (white)
            slide_title = section_title if chunk_idx == 0 else f"{section_title} (continued)"
            title_tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(0.65))
            title_tb.text_frame.word_wrap = True
            tp = title_tb.text_frame.paragraphs[0]
            tp.text = slide_title
            tp.font.bold = True
            tp.font.size = Pt(20)
            tp.font.color.rgb = white

            # Logo top-right
            add_logo(slide, 11.6, 0.22, width=1.15)

            # Pink accent line under bar (thin strip)
            acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.92), prs.slide_width, Inches(0.028))
            acc.fill.solid()
            acc.fill.fore_color.rgb = pink
            acc.line.fill.background()

            # Body: condensed bullets
            body_left = Inches(0.55)
            body_top = Inches(1.35)
            body_width = Inches(12.2)
            body_height = Inches(5.6)
            body_box = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
            body_box.text_frame.word_wrap = True
            body_box.text_frame.auto_size = None

            if bullets:
                for i, bullet_text in enumerate(bullets):
                    if i == 0:
                        para = body_box.text_frame.paragraphs[0]
                    else:
                        para = body_box.text_frame.add_paragraph()
                    para.text = bullet_text
                    para.font.size = Pt(12)
                    para.font.color.rgb = body_gray
                    para.space_before = Pt(6)
                    para.level = 0
            else:
                body_box.text_frame.paragraphs[0].text = "No additional details for this section."
                body_box.text_frame.paragraphs[0].font.size = Pt(12)
                body_box.text_frame.paragraphs[0].font.color.rgb = body_gray

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()
